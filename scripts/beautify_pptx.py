# One-off script to beautify docs/working-prompts/無人店面.pptx.
# Not part of the application; safe to delete after use.
import copy
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

SRC = r"docs/working-prompts/無人店面.pptx"

GOLD = RGBColor(0xC9, 0x97, 0x1F)      # Lucky gold accent (bold lead-in labels, table header, footer)
RED = RGBColor(0x8B, 0x1E, 0x1E)       # Lucky red accent (title-slide subtitle)
DARK = RGBColor(0x2B, 0x21, 0x18)      # warm dark charcoal (titles)
CREAM = RGBColor(0xFD, 0xF6, 0xE3)     # soft banding for table rows
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(SRC)


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_id_element = slides[index]
    rId = slide_id_element.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    xml_slides.remove(slide_id_element)
    # Also drop the now-unused relationship (and its target part, since
    # drop_rel removes the part too once nothing else references it) —
    # otherwise the orphaned part/relationship lingers and can collide with
    # part-name allocation on a later save (reproducible duplicate
    # 'ppt/slides/slideN.xml' zip-entry bug).
    prs.part.drop_rel(rId)


# --- 1. Remove the blank "Title Lorem Ipsum" leftover slide (index 1) ---
blank_idx = None
for i, slide in enumerate(prs.slides):
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame and title_shape.text_frame.text.strip() == "Title Lorem Ipsum":
        blank_idx = i
        break
if blank_idx is not None:
    delete_slide(prs, blank_idx)
    print(f"Removed blank slide at index {blank_idx}")
else:
    print("Blank slide not found (already removed?)")

# --- 2. Remove the empty ghost table on the finance slide ---
for slide in prs.slides:
    tables = [sh for sh in slide.shapes if getattr(sh, "has_table", False)]
    if len(tables) > 1:
        for t in tables:
            tbl = t.table
            all_empty = all(c.text.strip() == "" for row in tbl.rows for c in row.cells)
            if all_empty:
                t._element.getparent().remove(t._element)
                print(f"Removed empty ghost table {t.name!r} on a slide with {len(tables)} tables")

prs.save(SRC)
print("Saved after structural cleanup.")
