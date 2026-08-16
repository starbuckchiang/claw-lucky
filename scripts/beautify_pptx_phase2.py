# Phase 2: apply consistent color/typography polish to 無人店面.pptx.
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN

SRC = r"docs/working-prompts/無人店面.pptx"

GOLD = RGBColor(0xC9, 0x97, 0x1F)
RED = RGBColor(0x8B, 0x1E, 0x1E)
DARK = RGBColor(0x2B, 0x21, 0x18)
CREAM = RGBColor(0xFD, 0xF6, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(SRC)
slides = list(prs.slides)

# --- Slide 0 (title slide): fix the harsh pure-blue subtitle to a warm
# "Lucky Red" that reads clearly against the light khaki/beige band of the
# background photo (sampled avg ~ RGB(211,194,148)). ---
slide0 = slides[0]
for shape in slide0.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.type is not None:
                    try:
                        if run.font.color.rgb == RGBColor(0x00, 0x00, 0xFF):
                            run.font.color.rgb = RED
                            print("Fixed title-slide subtitle color -> RED:", repr(run.text))
                    except Exception:
                        pass

# --- Content slides (index 1..9): explicit title color + gold bold lead-ins ---
for i, slide in enumerate(slides):
    if i == 0 or i == len(slides) - 1:
        continue  # title slide + closing/contact slide styled separately
    title = slide.shapes.title
    if title is not None and title.has_text_frame:
        for para in title.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = DARK

    for shape in slide.shapes:
        if shape is title:
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.bold:
                        run.font.color.rgb = GOLD

print("Applied title/lead-in color pass.")
prs.save(SRC)
print("Saved.")
