# Phase 3: style the finance table + add a subtle branded footer to every
# content slide's already-present (currently empty) date placeholder.
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

SRC = r"docs/working-prompts/無人店面.pptx"

GOLD = RGBColor(0xC9, 0x97, 0x1F)
DARK = RGBColor(0x2B, 0x21, 0x18)
CREAM = RGBColor(0xFD, 0xF6, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(SRC)
slides = list(prs.slides)

# --- Style the finance table (header row + banded data rows + right-align
# the numeric "估算數值" column) ---
for slide in slides:
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        header = table.rows[0]
        for cell in header.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
        for ri, row in enumerate(table.rows):
            if ri == 0:
                continue
            band = CREAM if ri % 2 == 1 else WHITE
            for ci, cell in enumerate(row.cells):
                cell.fill.solid()
                cell.fill.fore_color.rgb = band
                for para in cell.text_frame.paragraphs:
                    if ci == 1:  # 估算數值 column -> right align for readability
                        para.alignment = PP_ALIGN.RIGHT
                    for run in para.runs:
                        if ci == 0:
                            run.font.bold = True
                            run.font.color.rgb = DARK
        print(f"Styled table {shape.name!r} ({len(table.rows)}x{len(table.columns)})")

# --- Fill every content slide's empty DATE placeholder with a small,
# consistent brand footer (purely additive — these placeholders already
# exist on every layout instance and were empty, so this is a safe,
# low-risk polish touch) ---
FOOTER_TEXT = "爪爪好運 Claw-Lucky ｜ 機密提案，僅供合作洽談使用"
styled_footers = 0
for i, slide in enumerate(slides):
    if i == 0:
        continue  # title slide already has its own footer-like text box
    for shape in slide.shapes:
        if shape.has_text_frame and "日期" in shape.name:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = FOOTER_TEXT
            run.font.size = Pt(10)
            run.font.color.rgb = GOLD
            run.font.italic = True
            styled_footers += 1

print(f"Added footer branding to {styled_footers} slides.")

prs.save(SRC)
print("Saved.")
